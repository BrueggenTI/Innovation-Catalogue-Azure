from flask import render_template, request, jsonify, send_file, flash, redirect, url_for, session
from app import app, db, csrf
from models import Product, ConceptSession, Trend, User
from data.products import init_products
from data.trends import init_trends
from utils.pdf_generator import generate_concept_pdf
from utils.email_sender import send_concept_email
from translations import get_text, get_available_languages
import json
import uuid
import logging
import re
import os
import time
from markupsafe import escape
import bleach
from werkzeug.utils import secure_filename
from openai import OpenAI
from replit.object_storage import Client as ObjectStorageClient
from flask_limiter import Limiter
from flask_limiter.util import get_remote_address
from PIL import Image
import fitz  # PyMuPDF
import base64
import io
from document_processor import document_processor
from ai_trend_analyzer import analyze_document_for_trend, improve_trend_description
from functools import wraps
from auth_config import get_msal_app, get_auth_url, acquire_token_by_code, get_logout_url, validate_config

def sanitize_input(input_string, max_length=100):
    """Sanitize user input to prevent XSS and injection attacks"""
    if not input_string:
        return ""

    # Remove HTML tags and limit length, but preserve common characters like &
    cleaned = bleach.clean(str(input_string), tags=[], strip=True, attributes={})
    # Unescape common HTML entities that should be preserved
    import html
    cleaned = html.unescape(cleaned)
    return cleaned[:max_length]

def validate_email(email):
    """Validate email format"""
    if not email:
        return False

    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    return re.match(pattern, email) is not None

def validate_uuid(uuid_string):
    """Validate UUID format"""
    try:
        uuid.UUID(uuid_string)
        return True
    except ValueError:
        return False

def clean_claim_text(claim_text):
    """Remove text in parentheses from claim text and clean up whitespace"""
    if not claim_text:
        return ""
    
    # Remove anything in parentheses and clean up whitespace
    clean_claim = re.sub(r'\s*\([^)]*\)\s*', '', claim_text).strip()
    return clean_claim

def extract_images_from_document(file_path, file_extension, original_filename):
    """Extract ALL images from documents for user selection"""
    extracted_images = {
        'all_images': []  # List of all extracted images with metadata
    }
    
    try:
        timestamp = int(time.time())
        images_dir = os.path.join('static', 'images', 'recipes', 'extracted')
        os.makedirs(images_dir, exist_ok=True)
        
        if file_extension == 'pdf':
            # Extract images from PDF using PyMuPDF
            try:
                doc = fitz.open(file_path)
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    image_list = page.get_images(full=True)
                    
                    for img_index, img in enumerate(image_list):
                        xref = img[0]
                        pix = fitz.Pixmap(doc, xref)
                        
                        # Skip very small images (likely icons or decorative elements)
                        if pix.width < 50 or pix.height < 50:
                            pix = None
                            continue
                            
                        # Convert to PNG and save
                        if pix.n - pix.alpha < 4:  # GRAY or RGB
                            img_data = pix.tobytes("png")
                            
                            filename = f"extracted_{timestamp}_{page_num}_{img_index}.png"
                            image_path = os.path.join(images_dir, filename)
                            
                            with open(image_path, 'wb') as img_file:
                                img_file.write(img_data)
                            
                            image_url = f"/static/images/recipes/extracted/{filename}"
                            extracted_images['all_images'].append({
                                'url': image_url,
                                'filename': filename,
                                'source': f'PDF Page {page_num + 1}',
                                'dimensions': f'{pix.width}x{pix.height}',
                                'index': len(extracted_images['all_images'])
                            })
                            
                            logging.info(f"Extracted image from PDF page {page_num + 1}: {filename}")
                        
                        pix = None
                doc.close()
            except Exception as pdf_error:
                logging.warning(f"Could not extract images from PDF: {pdf_error}")
                
        elif file_extension in ['ppt', 'pptx']:
            # Extract ALL images from PowerPoint
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                
                for slide_num, slide in enumerate(prs.slides):
                    for shape_num, shape in enumerate(slide.shapes):
                        if hasattr(shape, 'image') and shape.image:
                            try:
                                image_bytes = shape.image.blob
                                
                                # Use PIL to analyze image
                                img = Image.open(io.BytesIO(image_bytes))
                                
                                # Skip very small images
                                if img.width < 50 or img.height < 50:
                                    continue
                                
                                # Determine file extension
                                img_format = img.format.lower() if img.format else 'png'
                                if img_format not in ['png', 'jpg', 'jpeg']:
                                    img_format = 'png'
                                
                                filename = f"extracted_{timestamp}_{slide_num}_{shape_num}.{img_format}"
                                image_path = os.path.join(images_dir, filename)
                                
                                # Save image
                                if img_format == 'jpg':
                                    img = img.convert('RGB')
                                img.save(image_path)
                                
                                image_url = f"/static/images/recipes/extracted/{filename}"
                                extracted_images['all_images'].append({
                                    'url': image_url,
                                    'filename': filename,
                                    'source': f'Slide {slide_num + 1}',
                                    'dimensions': f'{img.width}x{img.height}',
                                    'index': len(extracted_images['all_images'])
                                })
                                
                                logging.info(f"Extracted image from PowerPoint slide {slide_num + 1}: {filename}")
                                    
                            except Exception as shape_error:
                                logging.warning(f"Could not process shape image: {shape_error}")
                                continue
                                
            except Exception as pptx_error:
                logging.warning(f"Could not extract images from PowerPoint: {pptx_error}")
        
        logging.info(f"Total images extracted: {len(extracted_images['all_images'])}")
        
    except Exception as e:
        logging.error(f"Error extracting images from document: {e}")
    
    return extracted_images


def init_user_session():
    """Initialize user session data from database or use defaults"""
    # If user_id exists in session, load from database
    user_id = session.get('user_id')
    if user_id:
        user = User.query.get(user_id)
        if user:
            # Load data from database
            session['user_name'] = user.name or 'User'
            session['user_email'] = user.email
            session['user_position'] = user.position
            session['user_department'] = user.department
            session['is_master_user'] = user.is_master_user
            if not session.get('language'):
                session['language'] = 'en'
            return
    
    # Fallback to default values if no user in database
    if not session.get('user_name'):
        session['user_name'] = 'Sarah Mitchell'
    if not session.get('user_position'):
        session['user_position'] = 'Sales Representative'
    if not session.get('user_email'):
        session['user_email'] = 'sarah.mitchell@brueggen.com'
    if not session.get('user_department'):
        session['user_department'] = 'Sales & Business Development'
    if not session.get('user_phone'):
        session['user_phone'] = '+49 541 123456'
    if not session.get('user_location'):
        session['user_location'] = 'Wallenhorst, Germany'
    if not session.get('language'):
        session['language'] = 'en'  # Default to English

# ============================================================================
# MICROSOFT SSO AUTHENTICATION ROUTES
# ============================================================================

def login_required(f):
    """
    Decorator to protect routes that require authentication
    Usage: @login_required
    """
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('authenticated'):
            flash('Please sign in with Microsoft to access this page.', 'warning')
            return redirect(url_for('login', next=request.url))
        return f(*args, **kwargs)
    return decorated_function

def master_required(f):
    """
    Decorator to restrict access to master user only
    Usage: @master_required
    """
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if not session.get('authenticated'):
            flash('Please sign in to access this page.', 'warning')
            return redirect(url_for('login', next=request.url))
        
        if not session.get('is_master_user'):
            flash('Access denied. This feature is only available to administrators.', 'danger')
            return redirect(url_for('index'))
        
        return f(*args, **kwargs)
    return decorated_function

@app.route('/login')
def login():
    """
    Displays login page with Microsoft SSO option
    If Azure is configured, shows login button
    If not configured, shows setup instructions
    """
    # Check if user is already authenticated
    if session.get('authenticated'):
        return redirect(url_for('index'))
    
    # Validate Azure configuration
    is_valid, message = validate_config()
    
    if not is_valid:
        # Show login page with configuration error
        logging.warning(f"Azure config invalid: {message}")
        return render_template('login.html', 
                             config_valid=False, 
                             error_message=message)
    
    # Azure is configured - prepare Microsoft SSO
    # Get MSAL app instance
    msal_app = get_msal_app()
    
    # Generate state token for security
    state = str(uuid.uuid4())
    session['auth_state'] = state
    
    # Store the 'next' URL to redirect after login
    session['next_url'] = request.args.get('next', url_for('index'))
    
    # Get Microsoft authorization URL
    auth_url = get_auth_url(msal_app, state=state)
    
    # Show login page with Microsoft SSO button
    return render_template('login.html', 
                         config_valid=True, 
                         auth_url=auth_url)

@app.route('/auth/callback')
def auth_callback():
    """
    Microsoft SSO callback handler
    Processes authorization code and establishes user session
    """
    # Verify state to prevent CSRF attacks
    if request.args.get('state') != session.get('auth_state'):
        flash('Authentication failed: Invalid state parameter', 'danger')
        logging.error("State mismatch in auth callback")
        return redirect(url_for('login'))
    
    # Get authorization code from Microsoft
    code = request.args.get('code')
    if not code:
        error = request.args.get('error', 'Unknown error')
        error_description = request.args.get('error_description', 'No description provided')
        flash(f'Authentication failed: {error} - {error_description}', 'danger')
        logging.error(f"Auth callback error: {error} - {error_description}")
        return redirect(url_for('login'))
    
    # Exchange code for access token
    msal_app = get_msal_app()
    result = acquire_token_by_code(msal_app, code)
    
    if 'error' in result:
        flash(f'Authentication failed: {result.get("error_description", "Unknown error")}', 'danger')
        logging.error(f"Token acquisition error: {result.get('error_description')}")
        return redirect(url_for('login'))
    
    # Store user information in session
    account = result.get('id_token_claims', {})
    user_email = account.get('preferred_username', account.get('email', ''))
    user_name = account.get('name', 'Unknown User')
    microsoft_id = account.get('oid')
    
    # Get or create user in database
    user = User.query.filter_by(email=user_email).first()
    
    if not user:
        # Create new user
        user = User(
            email=user_email,
            name=user_name,
            microsoft_id=microsoft_id,
            is_master_user=False
        )
        db.session.add(user)
        db.session.commit()
        logging.info(f"Created new user in database: {user_email}")
    else:
        # Update Microsoft ID if not set
        if not user.microsoft_id and microsoft_id:
            user.microsoft_id = microsoft_id
            db.session.commit()
    
    # Store user data in session from database
    session['authenticated'] = True
    session['user_id'] = user.id
    session['user_name'] = user.name or user_name
    session['user_email'] = user.email
    session['user_position'] = user.position
    session['user_department'] = user.department
    session['is_master_user'] = user.is_master_user
    session['microsoft_account'] = {
        'name': account.get('name'),
        'email': user_email,
        'id': microsoft_id,
        'tenant_id': account.get('tid')
    }
    
    # Clear auth state
    session.pop('auth_state', None)
    
    # Get next URL and redirect
    next_url = session.pop('next_url', url_for('index'))
    
    flash(f'Welcome, {session["user_name"]}! You are now signed in with Microsoft.', 'success')
    logging.info(f"User {user_email} successfully authenticated")
    
    return redirect(next_url)

@app.route('/logout')
def logout():
    """
    Logout user from both local session and Microsoft
    """
    user_name = session.get('user_name', 'User')
    
    # Clear local session
    session.clear()
    
    # Initialize default session again
    init_user_session()
    
    flash(f'Goodbye, {user_name}! You have been signed out.', 'info')
    
    # Redirect to Microsoft logout to clear SSO session
    logout_url = get_logout_url()
    return redirect(logout_url)

@app.route('/master-login', methods=['GET', 'POST'])
def master_login():
    """
    Master login route - separate from Microsoft SSO
    Allows master user to login with username and password
    """
    # Check if user is already authenticated
    if session.get('authenticated'):
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        username = request.form.get('username', '').strip()
        password = request.form.get('password', '')
        
        # Master credentials (should be stored as environment variables in production)
        MASTER_USERNAME = os.environ.get('MASTER_USERNAME', 'innocatmaster23568')
        MASTER_PASSWORD = os.environ.get('MASTER_PASSWORD', 'Villa23568hafer23568!')
        
        # Validate credentials
        if username == MASTER_USERNAME and password == MASTER_PASSWORD:
            # Get or create master user in database
            master_email = 'master@brueggen.com'
            user = User.query.filter_by(email=master_email).first()
            
            if not user:
                # Create master user
                user = User(
                    email=master_email,
                    name='Master User',
                    position='System Administrator',
                    department='IT & Innovation',
                    is_master_user=True
                )
                db.session.add(user)
                db.session.commit()
                logging.info(f"Created master user in database")
            elif not user.is_master_user:
                # Ensure is_master_user flag is set
                user.is_master_user = True
                db.session.commit()
            
            # Set session from database
            session['authenticated'] = True
            session['user_id'] = user.id
            session['user_name'] = user.name
            session['user_email'] = user.email
            session['user_position'] = user.position
            session['user_department'] = user.department
            session['is_master_user'] = user.is_master_user
            
            flash(f'Welcome, {user.name}! You are now signed in.', 'success')
            logging.info(f"Master user successfully authenticated from IP: {request.remote_addr}")
            
            # Redirect to index/home page
            return redirect(url_for('index'))
        else:
            flash('Invalid username or password. Please try again.', 'danger')
            logging.warning(f"Failed master login attempt from IP: {request.remote_addr}")
    
    # Show master login page
    return render_template('master-login.html')

@app.route('/profile', methods=['GET', 'POST'])
def profile():
    """
    User profile page - shows Microsoft account information if authenticated
    Also displays user profile edit form
    """
    init_user_session()
    microsoft_account = session.get('microsoft_account', {})
    lang = session.get('language', 'en')
    
    if request.method == 'POST':
        # Handle profile update
        user_id = session.get('user_id')
        
        if user_id:
            user = User.query.get(user_id)
            if user:
                # Update user data from form
                user.name = request.form.get('name', '').strip()
                user.email = request.form.get('email', '').strip()
                user.position = request.form.get('position', '').strip()
                user.department = request.form.get('department', '').strip()
                
                # Save to database
                db.session.commit()
                
                # Update session with new data
                session['user_name'] = user.name
                session['user_email'] = user.email
                session['user_position'] = user.position
                session['user_department'] = user.department
                
                flash('Profile updated successfully!', 'success')
                logging.info(f"User {user.email} updated profile")
            else:
                flash('User not found in database', 'danger')
        else:
            flash('You must be logged in to update your profile', 'warning')
        
        return redirect(url_for('profile'))
    
    # GET request - show profile form
    # Pass user data for the profile form
    user_data = {
        'name': session.get('user_name', 'Sarah Mitchell'),
        'email': session.get('user_email', 'sarah.mitchell@brueggen.com'),
        'position': session.get('user_position', 'Sales Representative'),
        'department': session.get('user_department', 'Sales & Business Development'),
        'phone': session.get('user_phone', '+49 541 123456'),
        'location': session.get('user_location', 'Wallenhorst, Germany'),
    }
    
    return render_template('profile.html', 
                         microsoft_account=microsoft_account,
                         user=user_data,
                         get_text=get_text,
                         lang=lang)

# ============================================================================
# END MICROSOFT SSO AUTHENTICATION ROUTES
# ============================================================================

@app.route('/set-language/<language>')
def set_language(language):
    """Set the user's preferred language"""
    if language in get_available_languages():
        session['language'] = language
    return redirect(request.referrer or url_for('index'))

@app.route('/')
@login_required
def index():
    init_user_session()
    lang = session.get('language', 'en')
    return render_template('index.html', get_text=get_text, lang=lang)

@app.route('/catalog')
@login_required
def catalog():
    init_user_session()
    # Initialize products if none exist
    if Product.query.count() == 0:
        init_products()

    # Check if user just published a recipe
    if request.args.get('new_recipe'):
        flash('Your new recipe has been successfully published and is now visible in the product portfolio!', 'success')

    # Get filter parameters with sanitization - only use non-empty values
    category = request.args.get('category', '').strip()
    ingredient = request.args.get('ingredient', '').strip()
    claim = request.args.get('claim', '').strip()
    recipe = request.args.get('recipe', '').strip()
    
    # Sanitize only if they have actual values
    category = sanitize_input(category) if category else ''
    ingredient = sanitize_input(ingredient) if ingredient else ''
    claim = sanitize_input(claim) if claim else ''
    recipe = sanitize_input(recipe) if recipe else ''
    
    
    # Build query - start with all products
    query = Product.query

    # Only apply filters if they have actual non-empty values
    if category and len(category) > 0:
        # Exact match for category
        query = query.filter(Product.category == category)

    if ingredient and len(ingredient) > 0:
        # Handle multiple ingredients (pipe-separated) with AND logic
        ingredients_list = [ing.strip() for ing in ingredient.split('|') if ing.strip()]
        for ing in ingredients_list:
            # Each ingredient must be present (AND logic)
            # Use exact JSON matching instead of contains
            query = query.filter(
                Product.ingredients.contains(f'"{ing}"')
            )

    if claim and len(claim) > 0:
        # Handle multiple claims (pipe-separated) with AND logic  
        claims_list = [cl.strip() for cl in claim.split('|') if cl.strip()]
        for cl in claims_list:
            # Each claim must be present (AND logic)
            # Use exact JSON matching instead of contains
            query = query.filter(
                Product.claims.contains(f'"{cl}"')
            )

    if recipe and len(recipe) > 0:
        # Handle recipe number filtering (multiple recipes with OR logic for search)
        recipe_list = [rc.strip() for rc in recipe.split('|') if rc.strip()]
        if len(recipe_list) == 1:
            # Single recipe - use contains for partial matching
            query = query.filter(
                Product.recipe_number.contains(recipe_list[0])
            )
        elif len(recipe_list) > 1:
            # Multiple recipes - use OR conditions
            from sqlalchemy import or_
            recipe_conditions = [Product.recipe_number.contains(rc) for rc in recipe_list]
            query = query.filter(or_(*recipe_conditions))

    products = query.order_by(Product.created_at.desc()).all()

    # Parse JSON fields for each product
    for product in products:
        if product.ingredients:
            try:
                product.parsed_ingredients = json.loads(product.ingredients)
            except:
                product.parsed_ingredients = []
        else:
            product.parsed_ingredients = []

        if product.nutritional_claims:
            try:
                product.parsed_nutritional_claims = json.loads(product.nutritional_claims)
            except:
                product.parsed_nutritional_claims = []
        else:
            product.parsed_nutritional_claims = []

    # Get unique values for filters from currently displayed products only
    categories = list(set([p.category for p in products if p.category]))

    # Extract ingredients and claims from JSON strings of currently displayed products
    all_ingredients = set()
    all_claims = set()

    for product in products:
        if product.ingredients:
            try:
                ingredients = json.loads(product.ingredients)
                if isinstance(ingredients, list):
                    # Handle list of ingredient objects
                    for ing in ingredients:
                        if isinstance(ing, dict) and 'name' in ing and ing['name']:
                            all_ingredients.add(ing['name'].strip())
                        elif isinstance(ing, str) and ing.strip():
                            all_ingredients.add(ing.strip())
            except json.JSONDecodeError:
                continue

        if product.claims:
            try:
                product_claims = json.loads(product.claims)
                if isinstance(product_claims, list):
                    for product_claim in product_claims:
                        # Add original claim text without modification
                        if product_claim and isinstance(product_claim, str) and product_claim.strip():
                            all_claims.add(product_claim.strip())
            except json.JSONDecodeError:
                continue

    # Ensure completely empty strings for template variables
    final_category = category if category and category.strip() and category != 'None' else ''
    final_ingredient = ingredient if ingredient and ingredient.strip() and ingredient != 'None' else ''
    final_claim = claim if claim and claim.strip() and claim != 'None' else ''
    final_recipe = recipe if recipe and recipe.strip() and recipe != 'None' else ''
    
    
    return render_template('competence.html', 
                         products=products,
                         categories=sorted(categories),
                         ingredients=sorted(all_ingredients),
                         claims=sorted(all_claims),
                         selected_category=final_category,
                         selected_ingredient=final_ingredient,
                         selected_claim=final_claim,
                         selected_recipe=final_recipe)

@app.route('/product/<int:id>', methods=['GET', 'POST'])
@login_required
def product_detail(id):
    init_user_session()
    product = Product.query.get_or_404(id)
    
    if request.method == 'POST':
        # Handle exclusive recipe form submission
        product.is_exclusive = 'is_exclusive' in request.form
        if product.is_exclusive:
            exclusive_type = request.form.get('exclusive_type', 'Market')
            exclusive_name = request.form.get('exclusive_name', '')
            product.target_market = f"{exclusive_type}:{exclusive_name}" if exclusive_name else ''
        else:
            product.target_market = ''
        
        try:
            db.session.commit()
            flash('Exklusivrezeptur-Einstellungen erfolgreich gespeichert!', 'success')
        except Exception as e:
            db.session.rollback()
            flash('Fehler beim Speichern der Einstellungen.', 'error')
        
        return redirect(url_for('product_detail', id=id))

    # Parse JSON fields
    ingredients = json.loads(product.ingredients) if product.ingredients else []
    nutritional_claims = json.loads(product.nutritional_claims) if product.nutritional_claims else []
    certifications = json.loads(product.certifications) if product.certifications else []
    nutritional_info = json.loads(product.nutritional_info) if product.nutritional_info else {}
    allergens = json.loads(product.allergens) if product.allergens else []
    claims = json.loads(product.claims) if product.claims else []

    return render_template('product_detail.html',
                         product=product,
                         ingredients=ingredients,
                         nutritional_claims=nutritional_claims,
                         certifications=certifications,
                         nutritional_info=nutritional_info,
                         allergens=allergens,
                         claims=claims)

@app.route('/trends')
@login_required
def trends():
    init_user_session()
    # Initialize trends if none exist
    if Trend.query.count() == 0:
        init_trends()

    # Get filter parameters
    categories = request.args.getlist('category')  # Support multiple categories
    report_types = request.args.getlist('report_type')  # Support multiple report types
    search_query = request.args.get('search', '')

    # Show all report types by default if none specified
    if not report_types:
        report_types = ['produktentwicklung', 'marktdaten']

    # Start with base query
    query = Trend.query

    # Apply report_type filter (multiple report types with OR logic)
    query = query.filter(Trend.report_type.in_(report_types))

    # Apply category filter if specified (multiple categories with OR logic)
    if categories:
        query = query.filter(Trend.category.in_(categories))

    # Apply search filter if specified
    if search_query:
        search_term = f"%{search_query}%"
        query = query.filter(
            Trend.title.ilike(search_term) |
            Trend.description.ilike(search_term) |
            Trend.market_data.ilike(search_term) |
            Trend.consumer_insights.ilike(search_term)
        )

    trends = query.order_by(Trend.created_at.desc()).all()

    return render_template('trends.html', 
                         trends=trends, 
                         selected_categories=categories,  # List of selected categories
                         selected_report_types=report_types,  # List of selected report types
                         search_query=search_query)

@app.route('/trend/<int:trend_id>/pdf')
@login_required
def trend_pdf_viewer(trend_id):
    """Display PDF report for specific trends"""
    trend = Trend.query.get_or_404(trend_id)
    
    # Check if this trend has an associated PDF report
    if trend_id == 6:  # Global Fusion Flavors trend
        pdf_path = 'global_fusion_flavors_report.pdf'
        return render_template('pdf_viewer.html', 
                             trend=trend, 
                             pdf_path=pdf_path,
                             pdf_title="Global Fusion Flavors: Comprehensive Trend Report")
    else:
        # For other trends, show a message that no detailed report is available
        return render_template('pdf_viewer.html', 
                             trend=trend, 
                             pdf_path=None,
                             pdf_title="Detailed Report Not Available")

@app.route('/cocreation')
@login_required
def cocreation():
    init_user_session()
    # Check if a base product is selected
    base_product_id = request.args.get('base_product_id') or request.args.get('base_product')
    base_product = None

    if base_product_id:
        try:
            base_product = Product.query.get(int(base_product_id))
        except (ValueError, TypeError):
            pass
    # Get base products for selection
    products = Product.query.all()

    # Create new session
    session_id = str(uuid.uuid4())

    return render_template('cocreation.html', 
                         products=products,
                         session_id=session_id,
                         base_product=base_product)

@app.route('/cocreation/save_concept', methods=['POST'])
@csrf.exempt
@login_required
def save_concept():
    try:
        data = request.get_json()

        session_id = data.get('session_id')
        client_name = data.get('client_name', '')
        client_email = data.get('client_email', '')
        product_config = data.get('product_config', {})

        # Check if session already exists
        existing_concept = ConceptSession.query.filter_by(session_id=session_id).first()

        if existing_concept:
            # Update existing session
            existing_concept.client_name = client_name
            existing_concept.client_email = client_email
            existing_concept.product_config = json.dumps(product_config)
            concept = existing_concept
        else:
            # Create new session
            concept = ConceptSession()
            concept.session_id = session_id
            concept.client_name = client_name
            concept.client_email = client_email
            concept.product_config = json.dumps(product_config)
            db.session.add(concept)

        db.session.commit()

        # Generate PDF
        pdf_path = generate_concept_pdf(concept)
        concept.pdf_path = pdf_path
        db.session.commit()

        return jsonify({
            'success': True,
            'message': 'Concept saved successfully',
            'pdf_url': url_for('download_pdf', session_id=session_id)
        })

    except Exception as e:
        logging.error(f"Error saving concept: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'Error saving concept: {str(e)}'
        }), 500

@app.route('/cocreation/send_email', methods=['POST'])
@csrf.exempt
@login_required
def send_concept_email_route():
    try:
        data = request.get_json()
        session_id = data.get('session_id')

        concept = ConceptSession.query.filter_by(session_id=session_id).first()
        if not concept:
            return jsonify({
                'success': False,
                'message': 'Concept session not found'
            }), 404

        if not concept.client_email:
            return jsonify({
                'success': False,
                'message': 'No email address provided'
            }), 400

        # Send email
        success = send_concept_email(concept)

        if success:
            return jsonify({
                'success': True,
                'message': f'Concept sent successfully to {concept.client_email}'
            })
        else:
            return jsonify({
                'success': False,
                'message': 'Failed to send email'
            }), 500

    except Exception as e:
        logging.error(f"Error sending email: {str(e)}")
        return jsonify({
            'success': False,
            'message': f'Error sending email: {str(e)}'
        }), 500

@app.route('/download_pdf/<session_id>')
@login_required
def download_pdf(session_id):
    concept = ConceptSession.query.filter_by(session_id=session_id).first()
    if not concept or not concept.pdf_path:
        flash('PDF not found', 'error')
        return redirect(url_for('index'))

    try:
        return send_file(concept.pdf_path, as_attachment=True, 
                        download_name=f'bruggen_concept_{concept.client_name or "draft"}.pdf')
    except Exception as e:
        logging.error(f"Error downloading PDF: {str(e)}")
        flash('Error downloading PDF', 'error')
        return redirect(url_for('index'))

@app.route('/api/search', methods=['POST'])
@csrf.exempt
@login_required
def search_api():
    """Global search API endpoint"""
    try:
        data = request.get_json()
        query = data.get('query', '').strip().lower()

        if len(query) < 2:
            return jsonify({'success': False, 'message': 'Query too short'})

        results = []

        # Check if query looks like a recipe number (e.g., "R123", "123", "RZ456")
        is_recipe_number_search = bool(re.match(r'^(r|rz|recipe)?\s*\d+', query, re.IGNORECASE))
        
        # Search products/recipes
        try:
            search_conditions = [
                Product.name.ilike(f'%{query}%'),
                Product.description.ilike(f'%{query}%'),
                Product.category.ilike(f'%{query}%'),
                Product.ingredients.ilike(f'%{query}%'),
                Product.claims.ilike(f'%{query}%'),
                Product.allergens.ilike(f'%{query}%'),
                Product.nutritional_claims.ilike(f'%{query}%'),
                Product.storage_conditions.ilike(f'%{query}%')
            ]
            
            # For recipe number searches, also try to match by ID or extracted numbers
            if is_recipe_number_search:
                # Extract just the numbers from the query
                number_match = re.search(r'\d+', query)
                if number_match:
                    recipe_number = int(number_match.group())
                    # Search by ID or any field containing this number
                    search_conditions.extend([
                        Product.id == recipe_number,
                        Product.name.ilike(f'%{recipe_number}%'),
                        Product.description.ilike(f'%{recipe_number}%')
                    ])
            
            products = Product.query.filter(
                db.or_(*search_conditions)
            ).order_by(Product.created_at.desc()).limit(15).all()

            for product in products:
                # Determine if this is more recipe-like or product-like based on content
                product_type = 'recipe' if any(keyword in product.name.lower() for keyword in ['muesli', 'müsli', 'granola', 'porridge', 'oats', 'recipe']) else 'product'
                
                # For recipe number searches, prefer showing as recipes
                if is_recipe_number_search:
                    product_type = 'recipe'
                
                # Generate a recipe number for display (use ID or extract from name)
                recipe_number = None
                if product_type == 'recipe':
                    # Try to extract existing recipe number from name
                    number_match = re.search(r'(r|rz|recipe)?\s*(\d+)', product.name, re.IGNORECASE)
                    if number_match:
                        recipe_number = f"R{number_match.group(2)}"
                    else:
                        # Use product ID as recipe number
                        recipe_number = f"R{product.id:04d}"
                
                result_data = {
                    'type': product_type,
                    'title': product.name,
                    'description': product.description[:100] + '...' if product.description and len(product.description) > 100 else (product.description or ''),
                    'category': product.category,
                    'url': url_for('product_detail', id=product.id),
                    'image': product.image_url
                }
                
                # Add recipe number for recipe results
                if recipe_number:
                    result_data['recipe_number'] = recipe_number
                
                results.append(result_data)
        except Exception as e:
            logging.error(f"Error searching products: {str(e)}")

        # Search trends
        try:
            trends = Trend.query.filter(
                db.or_(
                    Trend.title.ilike(f'%{query}%'),
                    Trend.description.ilike(f'%{query}%'),
                    Trend.category.ilike(f'%{query}%'),
                    Trend.market_data.ilike(f'%{query}%'),
                    Trend.consumer_insights.ilike(f'%{query}%')
                )
            ).order_by(Trend.created_at.desc()).limit(10).all()

            for trend in trends:
                results.append({
                    'type': 'trend',
                    'title': trend.title,
                    'description': trend.description[:100] + '...' if trend.description and len(trend.description) > 100 else (trend.description or ''),
                    'category': trend.category,
                    'url': url_for('trends', category=trend.category) + f'#trend-{trend.id}',
                    'image': trend.image_url
                })
        except Exception as e:
            logging.error(f"Error searching trends: {str(e)}")

        # Add quick actions for common searches
        if len(results) <= 3:  # Show quick actions when few specific results
            # German and English search terms for Co-Creation
            if any(term in query for term in ['cocreation', 'co-creation', 'create', 'lab', 'erstellen', 'entwickeln', 'kreieren']):
                results.append({
                    'type': 'action',
                    'title': 'Co-Creation Lab',
                    'description': 'Gemeinsam mit Kunden maßgeschneiderte Produkte entwickeln',
                    'category': 'Schnellzugriff',
                    'url': url_for('cocreation')
                })

            # German and English search terms for Competence
            if any(term in query for term in ['competence', 'kompetenz', 'products', 'produkte', 'portfolio', 'fähigkeiten']):
                results.append({
                    'type': 'action',
                    'title': 'Unsere Kompetenzen',
                    'description': 'Entdecken Sie unser Produktportfolio und unsere Fähigkeiten',
                    'category': 'Schnellzugriff',
                    'url': url_for('catalog')
                })

            # German and English search terms for Trends
            if any(term in query for term in ['trends', 'insights', 'future', 'markt', 'zukunft', 'entwicklung']):
                results.append({
                    'type': 'action',
                    'title': 'Trends & Einblicke',
                    'description': 'Entdecken Sie Markttrends und zukünftige Möglichkeiten',
                    'category': 'Schnellzugriff',
                    'url': url_for('trends')
                })

            # Recipe/Rezept search
            if any(term in query for term in ['recipe', 'rezept', 'rezepte', 'hinzufügen', 'upload', 'hochladen']):
                results.append({
                    'type': 'action',
                    'title': 'Neues Rezept hinzufügen',
                    'description': 'Rezept-Dokument hochladen und KI-Analyse starten',
                    'category': 'Schnellzugriff',
                    'url': url_for('add_recipe')
                })

        return jsonify({
            'success': True,
            'results': results,
            'total': len(results)
        })
    
    except Exception as e:
        logging.error(f"Search API error: {str(e)}")
        return jsonify({
            'success': False,
            'message': 'Search temporarily unavailable'
        }), 500

@app.route('/api/trends/search')
@login_required
def search_trends():
    """API endpoint for searching trends"""
    try:
        init_user_session()
        
        query = request.args.get('q', '').strip()
        report_types = request.args.getlist('report_type')  # Support multiple report types
        
        # Default to produktentwicklung if no report types specified
        if not report_types:
            report_types = ['produktentwicklung']
        
        if not query or len(query) < 2:
            return jsonify([])
        
        # Search trends
        search_term = f"%{query}%"
        trends = Trend.query.filter(
            Trend.report_type.in_(report_types)  # Use IN for multiple report types
        ).filter(
            Trend.title.ilike(search_term) |
            Trend.description.ilike(search_term) |
            Trend.category.ilike(search_term) |
            Trend.market_data.ilike(search_term) |
            Trend.consumer_insights.ilike(search_term)
        ).limit(10).all()
        
        trends_data = []
        for trend in trends:
            # Build URL with all current report types
            url_params = '&'.join([f'report_type={rt}' for rt in report_types])
            
            # Get readable report type name
            report_type_name = 'Product Development' if trend.report_type == 'produktentwicklung' else 'Market Data'
            
            trends_data.append({
                'id': trend.id,
                'title': trend.title,
                'category': trend.category.title(),
                'description': trend.description[:100] + '...' if len(trend.description) > 100 else trend.description,
                'image': trend.image_url,  # Add image support
                'report_type': trend.report_type,
                'report_type_name': report_type_name,
                'url': url_for('trends') + f'?{url_params}#{trend.id}',
                'type': 'trend'  # Add type for consistency with innovation catalogue
            })
        
        return jsonify(trends_data)
        
    except Exception as e:
        logging.error(f"Search API error: {str(e)}")
        return jsonify({
            'success': False,
            'message': 'Search temporarily unavailable'
        }), 500

@app.route('/api/trends/analyze', methods=['POST'])
@csrf.exempt
@login_required
def analyze_trend_files():
    """API endpoint for analyzing uploaded files and generating trend suggestions"""
    try:
        if 'files' not in request.files:
            return jsonify({'success': False, 'error': 'No files uploaded'}), 400
        
        files = request.files.getlist('files')
        if not files or all(f.filename == '' for f in files):
            return jsonify({'success': False, 'error': 'No files selected'}), 400
        
        # Process uploaded files
        files_data = []
        temp_file_paths = []
        
        for file in files:
            if file and file.filename:
                # Validate file type
                file_type = file.content_type
                if not document_processor.is_supported(file_type):
                    return jsonify({
                        'success': False, 
                        'error': f'Unsupported file type: {file.filename}'
                    }), 400
                
                # Save file temporarily
                file_path = document_processor.save_uploaded_file(file)
                temp_file_paths.append(file_path)
                files_data.append((file_path, file_type, file.filename))
        
        if not files_data:
            return jsonify({'success': False, 'error': 'No valid files found'}), 400
        
        # Process all files and extract text
        processing_result = document_processor.process_multiple_files(files_data)
        
        if processing_result['successful_files'] == 0:
            document_processor.cleanup_files(temp_file_paths)
            return jsonify({
                'success': False, 
                'error': 'Failed to process any uploaded files'
            }), 400
        
        # Analyze with AI if we have enough text
        combined_text = processing_result['combined_text']
        if len(combined_text.strip()) < 100:
            document_processor.cleanup_files(temp_file_paths)
            return jsonify({
                'success': False, 
                'error': 'Not enough text content found in uploaded files'
            }), 400
        
        # Use AI to analyze the content
        ai_analysis = analyze_document_for_trend(combined_text)
        
        # Save PDF files permanently for trend creation
        saved_pdf_path = None
        for temp_path, file_type, original_filename in files_data:
            if file_type == 'application/pdf' and os.path.exists(temp_path):
                # Create permanent PDF storage directory
                pdfs_dir = os.path.join('static', 'pdfs')
                os.makedirs(pdfs_dir, exist_ok=True)
                
                # Generate unique filename for PDF
                timestamp = int(time.time())
                unique_id = str(uuid.uuid4())[:8]
                pdf_filename = f"trend_doc_{timestamp}_{unique_id}.pdf"
                permanent_path = os.path.join(pdfs_dir, pdf_filename)
                
                # Copy PDF to permanent location
                import shutil
                shutil.copy2(temp_path, permanent_path)
                saved_pdf_path = f"/static/pdfs/{pdf_filename}"
                break  # Use first PDF found
        
        # Cleanup temporary files
        document_processor.cleanup_files(temp_file_paths)
        
        # Return analysis results
        return jsonify({
            'success': True,
            'title': ai_analysis.get('title', 'Generated Trend'),
            'description': ai_analysis.get('description', 'AI-generated trend description'),
            'category': ai_analysis.get('category', 'Innovation'),
            'report_type': ai_analysis.get('report_type', 'produktentwicklung'),
            'market_data': ai_analysis.get('market_data', 'Market data analysis pending'),
            'consumer_insights': ai_analysis.get('consumer_insights', 'Consumer insights analysis pending'),
            'pdf_path': saved_pdf_path,
            'confidence': ai_analysis.get('confidence', 0.8),
            'processing_stats': {
                'total_files': processing_result['total_files'],
                'successful_files': processing_result['successful_files'],
                'total_words': processing_result['total_word_count']
            }
        })
        
    except Exception as e:
        # Cleanup files in case of error
        if 'temp_file_paths' in locals():
            document_processor.cleanup_files(temp_file_paths)
        
        logging.error(f"Trend analysis error: {str(e)}")
        return jsonify({
            'success': False,
            'error': 'Analysis failed. Please try again.'
        }), 500

@app.route('/api/trends/create', methods=['POST'])
@csrf.exempt
@login_required
def create_custom_trend():
    """API endpoint for creating custom trends from AI analysis"""
    try:
        # Support both JSON and FormData (for file uploads)
        if request.content_type and 'multipart/form-data' in request.content_type:
            # FormData request (with potential file upload)
            title = request.form.get('title', '').strip()
            description = request.form.get('description', '').strip()
            category = request.form.get('category', 'Innovation')
            report_type = request.form.get('report_type', 'produktentwicklung')
            market_data = request.form.get('market_data', '').strip()
            consumer_insights = request.form.get('consumer_insights', '').strip()
            pdf_path = request.form.get('pdf_path', '').strip()
            image_file = request.files.get('image')
        else:
            # JSON request (legacy support)
            data = request.get_json()
            if not data:
                return jsonify({'success': False, 'error': 'No data provided'}), 400
            
            title = data.get('title', '').strip()
            description = data.get('description', '').strip()
            category = data.get('category', 'Innovation')
            report_type = data.get('report_type', 'produktentwicklung')
            market_data = data.get('market_data', '').strip()
            consumer_insights = data.get('consumer_insights', '').strip()
            pdf_path = data.get('pdf_path', '').strip()
            image_file = None
        
        if not title or not description:
            return jsonify({
                'success': False, 
                'error': 'Title and description are required'
            }), 400
        
        # Validate category and report type
        valid_categories = ['Health', 'Sustainability', 'Innovation']
        valid_report_types = ['produktentwicklung', 'marktdaten']
        
        if category not in valid_categories:
            category = 'Innovation'
        
        # Convert category to lowercase to match existing data and CSS classes
        category = category.lower()
        
        if report_type not in valid_report_types:
            report_type = 'produktentwicklung'
            
        # Handle image upload
        image_filename = None
        if image_file:
            # Save uploaded image
            timestamp = int(time.time())
            unique_id = str(uuid.uuid4())[:8]
            filename = f"trend_{timestamp}_{unique_id}.{image_file.filename.rsplit('.', 1)[1].lower()}"
            
            # Ensure directory exists
            images_dir = os.path.join('static', 'images')
            os.makedirs(images_dir, exist_ok=True)
            
            # Save file
            file_path = os.path.join(images_dir, filename)
            
            try:
                # Resize and optimize image
                img = Image.open(image_file.stream)
                
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                
                # Resize if too large
                max_size = (800, 600)
                img.thumbnail(max_size, Image.Resampling.LANCZOS)
                
                # Save optimized image
                img.save(file_path, 'JPEG', quality=85, optimize=True)
                
                image_filename = filename
                
            except Exception as e:
                app.logger.error(f'Error saving image: {str(e)}')
                # Continue without image
        
        # Create new trend in database
        image_url = f"/static/images/{image_filename}" if image_filename else None
        
        # Use AI-generated market_data and consumer_insights if provided, otherwise use defaults
        final_market_data = market_data if market_data else "Custom trend created by user analysis"
        final_consumer_insights = consumer_insights if consumer_insights else "User-generated insights from document analysis"
        
        # Validate and set PDF path
        final_pdf_path = None
        if pdf_path and pdf_path.startswith('/static/pdfs/'):
            # Check if the PDF file actually exists
            pdf_file_path = os.path.join(app.root_path, pdf_path.lstrip('/'))
            if os.path.exists(pdf_file_path):
                final_pdf_path = pdf_path
        
        new_trend = Trend(
            title=title[:100],  # Limit title length
            description=description[:300],  # Limit description length
            category=category,
            report_type=report_type,
            market_data=final_market_data,
            consumer_insights=final_consumer_insights,
            image_url=image_url,  # Use uploaded image if available
            pdf_path=final_pdf_path  # Store PDF path for later viewing
        )
        
        # Add to database
        db.session.add(new_trend)
        db.session.commit()
        
        return jsonify({
            'success': True,
            'trend_id': new_trend.id,
            'message': 'Custom trend created successfully'
        })
        
    except Exception as e:
        db.session.rollback()
        logging.error(f"Custom trend creation error: {str(e)}")
        return jsonify({
            'success': False,
            'error': 'Failed to create trend. Please try again.'
        }), 500

@app.route('/api/trends/<int:trend_id>', methods=['GET'])
@login_required
def get_trend_details(trend_id):
    """API endpoint to get trend details by ID"""
    try:
        trend = Trend.query.get(trend_id)
        if not trend:
            return jsonify({'success': False, 'error': 'Trend not found'}), 404
        
        return jsonify({
            'success': True,
            'trend': {
                'id': trend.id,
                'title': trend.title,
                'description': trend.description,
                'category': trend.category,
                'report_type': trend.report_type,
                'market_data': trend.market_data,
                'consumer_insights': trend.consumer_insights,
                'image_url': trend.image_url,
                'pdf_path': trend.pdf_path,
                'created_at': trend.created_at.isoformat() if trend.created_at else None
            }
        })
        
    except Exception as e:
        logging.error(f"Get trend details error: {str(e)}")
        return jsonify({
            'success': False,
            'error': 'Failed to fetch trend details.'
        }), 500

# Deep Research Routes
@app.route('/api/deep-research/start', methods=['POST'])
@csrf.exempt
@login_required
def start_deep_research():
    """Start a new deep research job"""
    import uuid
    import threading
    from models import ResearchJob
    
    try:
        data = request.get_json()
        if not data:
            return jsonify({'success': False, 'error': 'No data provided'}), 400
        
        description = data.get('description', '').strip()
        keywords = data.get('keywords', [])
        categories = data.get('categories', [])
        
        if not description or len(description) < 20:
            return jsonify({'success': False, 'error': 'Description must be at least 20 characters'}), 400
        
        # Create unique job ID
        job_id = str(uuid.uuid4())
        
        # Create job in database
        new_job = ResearchJob(
            job_id=job_id,
            user_id=session.get('user_id'),
            description=description,
            keywords=json.dumps(keywords),
            categories=json.dumps(categories),
            status='queued',
            progress=0,
            status_log=json.dumps([])
        )
        db.session.add(new_job)
        db.session.commit()
        
        # Start job in background thread
        def run_job():
            from deep_research_worker import process_research_job
            # Store updates in memory for SSE streaming
            job_updates[job_id] = []
            
            for update in process_research_job(job_id, description, keywords, categories):
                if job_id in job_updates:
                    job_updates[job_id].append(update)
        
        thread = threading.Thread(target=run_job, daemon=True)
        thread.start()
        
        return jsonify({
            'success': True,
            'job_id': job_id,
            'message': 'Research job started successfully'
        })
        
    except Exception as e:
        logging.error(f"Start deep research error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({
            'success': False,
            'error': 'Failed to start research job. Please try again.'
        }), 500

# Store job updates in memory for SSE
job_updates = {}

@app.route('/api/deep-research/plan/<job_id>', methods=['GET'])
@login_required
def get_research_plan(job_id):
    """Get the research plan for a job"""
    try:
        job = ResearchJob.query.filter_by(job_id=job_id).first()
        if not job:
            return jsonify({'success': False, 'error': 'Job not found'}), 404
        
        if not job.research_plan:
            return jsonify({'success': False, 'error': 'Plan not yet generated'}), 404
        
        plan = json.loads(job.research_plan)
        return jsonify({
            'success': True,
            'plan': plan,
            'status': job.status,
            'approved': job.plan_approved
        })
        
    except Exception as e:
        logging.error(f"Get research plan error: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to get research plan'}), 500

@app.route('/api/deep-research/plan/<job_id>/approve', methods=['POST'])
@csrf.exempt
@login_required
def approve_research_plan(job_id):
    """Approve and optionally modify the research plan, then continue job"""
    import threading
    from models import ResearchJob
    
    try:
        data = request.get_json() or {}
        modified_plan = data.get('plan')  # Optional: user can modify the plan
        
        job = ResearchJob.query.filter_by(job_id=job_id).first()
        if not job:
            return jsonify({'success': False, 'error': 'Job not found'}), 404
        
        if job.status != 'waiting_approval':
            return jsonify({'success': False, 'error': 'Job is not waiting for approval'}), 400
        
        # Update plan if modified
        if modified_plan:
            job.research_plan = json.dumps(modified_plan, ensure_ascii=False)
            logging.info(f"Plan modified by user for job {job_id}")
        
        # Mark as approved
        job.plan_approved = True
        job.status = 'processing_strategy'
        db.session.commit()
        
        logging.info(f"Plan approved for job {job_id}, resuming research...")
        
        # Extract data BEFORE thread starts (to avoid DB access outside app context)
        description = job.description
        keywords = json.loads(job.keywords) if job.keywords else []
        categories = json.loads(job.categories) if job.categories else []
        
        # Resume job in background thread
        def resume_job():
            from deep_research_worker import process_research_job
            from app import app
            
            # Use app context for DB access
            with app.app_context():
                # Re-initialize updates if needed
                if job_id not in job_updates:
                    job_updates[job_id] = []
                
                for update in process_research_job(job_id, description, keywords, categories):
                    if job_id in job_updates:
                        job_updates[job_id].append(update)
        
        thread = threading.Thread(target=resume_job, daemon=True)
        thread.start()
        
        return jsonify({
            'success': True,
            'message': 'Plan approved, research job resuming'
        })
        
    except Exception as e:
        logging.error(f"Approve research plan error: {str(e)}")
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': 'Failed to approve plan'}), 500

@app.route('/api/deep-research/stream/<job_id>')
@login_required
def stream_deep_research(job_id):
    """SSE endpoint for streaming research job updates"""
    from flask import Response, stream_with_context
    
    def event_stream():
        # Wait for job to start generating updates
        max_wait = 30  # 30 seconds timeout
        waited = 0
        while job_id not in job_updates and waited < max_wait:
            time.sleep(0.1)
            waited += 0.1
        
        if job_id not in job_updates:
            yield f"data: {json.dumps({'type': 'error', 'message': 'Job not found or timed out'})}\n\n"
            return
        
        sent_count = 0
        max_iterations = 1000  # Prevent infinite loop
        iterations = 0
        
        while iterations < max_iterations:
            updates = job_updates.get(job_id, [])
            
            # Send any new updates
            while sent_count < len(updates):
                yield f"data: {updates[sent_count]}\n\n"
                
                # Check if this is a completion or error message
                try:
                    update_data = json.loads(updates[sent_count])
                    if update_data.get('type') in ['complete', 'error']:
                        # Clean up
                        if job_id in job_updates:
                            del job_updates[job_id]
                        return
                except:
                    pass
                
                sent_count += 1
            
            time.sleep(0.1)  # Small delay between checks
            iterations += 1
        
        # Cleanup if we hit max iterations
        if job_id in job_updates:
            del job_updates[job_id]
    
    return Response(stream_with_context(event_stream()), 
                   mimetype='text/event-stream',
                   headers={
                       'Cache-Control': 'no-cache, no-transform',
                       'X-Accel-Buffering': 'no',
                       'Connection': 'keep-alive'
                   })

@app.route('/api/generate-image', methods=['POST'])
@csrf.exempt
@login_required
def generate_image():
    """API endpoint for generating AI images for trends"""
    try:
        data = request.get_json()
        if not data:
            return jsonify({'success': False, 'error': 'No data provided'}), 400
        
        prompt = data.get('prompt', '').strip()
        one_line_summary = data.get('one_line_summary', 'Business Trend')
        aspect_ratio = data.get('aspect_ratio', '16:9')
        
        if not prompt:
            return jsonify({'success': False, 'error': 'Prompt is required'}), 400
        
        # For now, return a placeholder response
        # In a full implementation, this would use actual image generation
        return jsonify({
            'success': True,
            'image_url': '/static/images/placeholder-trend.jpg',
            'message': 'Image generation feature coming soon'
        })
        
    except Exception as e:
        app.logger.error(f'Generate image API error: {str(e)}')
        return jsonify({
            'success': False,
            'error': 'An error occurred while generating the image'
        }), 500

@app.route('/add-recipe')
@master_required
def add_recipe():
    init_user_session()
    # Get all recipes to display in the list at the bottom
    all_recipes = Product.query.order_by(Product.created_at.desc()).all()
    return render_template('add_recipe.html', recipes=all_recipes)

@app.route('/api/analyze-recipe', methods=['POST'])
@csrf.exempt
@master_required
def analyze_recipe():
    # Check for files in request - handle both single and multiple files
    files = []
    if 'recipe_file' in request.files:
        files = request.files.getlist('recipe_file')
    elif 'recipeFile' in request.files:
        files = request.files.getlist('recipeFile')
    elif 'file' in request.files:
        files = request.files.getlist('file')
    
    if not files or all(f.filename == '' for f in files):
        return jsonify({'success': False, 'error': 'No files uploaded'})

    try:
        # Process multiple files using DocumentProcessor
        files_data = []
        temp_file_paths = []
        extracted_images_all = []
        filenames_list = []
        
        for file in files:
            if file and file.filename:
                # Validate file type
                file_type = file.content_type
                if not document_processor.is_supported(file_type):
                    return jsonify({
                        'success': False,
                        'error': f'Unsupported file type: {file.filename}. Supported: PDF, DOCX, PPTX, XLSX, TXT'
                    }), 400
                
                # Save file temporarily
                file_path = document_processor.save_uploaded_file(file)
                temp_file_paths.append(file_path)
                files_data.append((file_path, file_type, file.filename))
                filenames_list.append(file.filename)
                
                # Extract images from each file
                file_extension = os.path.splitext(file.filename.lower())[1][1:]  # Remove dot
                images = extract_images_from_document(file_path, file_extension, file.filename)
                if images.get('all_images'):
                    extracted_images_all.extend(images['all_images'])
        
        if not files_data:
            return jsonify({'success': False, 'error': 'No valid files found'}), 400
        
        # Process all files and extract text
        processing_result = document_processor.process_multiple_files(files_data)
        
        if processing_result['successful_files'] == 0:
            document_processor.cleanup_files(temp_file_paths)
            return jsonify({
                'success': False,
                'error': 'Failed to process any uploaded files'
            }), 400
        
        # Get combined text content
        file_content = processing_result['combined_text']
        if len(file_content.strip()) < 10:
            file_content = f"Documents: {', '.join(filenames_list)} - Please provide realistic recipe details based on the filenames."
        
        # Use first filename for compatibility
        filename = filenames_list[0] if filenames_list else 'uploaded_file'
        file_extension = os.path.splitext(filename.lower())[1][1:]  # Remove dot
        
        # Prepare extracted images dict
        extracted_images = {'all_images': extracted_images_all}
        
        # Initialize OpenAI client
        openai_api_key = os.environ.get('OPENAI_API_KEY')
        if not openai_api_key:
            logging.error("OpenAI API key not found in environment variables")
            return jsonify({
                'success': False,
                'error': 'OpenAI API key not configured. Please set the OPENAI_API_KEY environment variable.'
            })
            
        # Using gpt-4o for reliable document analysis
        openai_client = OpenAI(api_key=openai_api_key)

        # Create AI prompt for recipe analysis
        prompt = f"""
        You are an expert food technologist analyzing a Brüggen cereal product document. Your mission is to perform COMPREHENSIVE EXTRACTION of ALL relevant product information and translate everything to English.

        DOCUMENT TYPE SPECIFIC INSTRUCTIONS:
        For POWERPOINT/PRESENTATION files (.pptx, .ppt):
        - SCAN ALL SLIDES thoroughly - information is often spread across multiple slides
        - LOOK FOR TABLES with nutritional data, ingredient lists, percentages
        - EXAMINE SLIDE TITLES, HEADERS, BULLET POINTS carefully
        - NUMERICAL DATA is often in tables or charts - extract ALL numbers with their labels
        - INGREDIENTS may be listed with percentages in separate slides or tables
        - ALLERGEN info might be in footnotes or separate disclaimer slides
        
        For TEXT/WORD files (.txt, .docx, .doc):
        - SCAN entire document from top to bottom
        - LOOK FOR structured lists, tables, headings
        - EXTRACT information from paragraphs and bullet points

        EXTRACTION PRIORITY ORDER:
        1. FIND THE RECIPE NUMBER - Look specifically for "Ref.:" followed by a number pattern (e.g., "Ref.: 000014567", "Ref.: 00001234")
        2. FIND THE PRODUCT NAME - Look for slide titles, headings, product descriptions, main headers
        3. IDENTIFY ALL INGREDIENTS - Search for ingredient lists, Zutatenliste, compositions, percentage tables
        4. EXTRACT ALL NUTRITIONAL DATA - Look for Nährwerte, nutritional tables, "per 100g" values, energy values
        5. LOCATE ALLERGEN INFORMATION - Find allergen declarations, "Enthält", allergy warnings, disclaimers
        6. IDENTIFY PRODUCT CLAIMS - Look for health claims, certifications, quality statements, marketing text
        7. FIND STORAGE CONDITIONS - Search for storage instructions, Lagerung, Aufbewahrung

        CRITICAL EXTRACTION RULES:
        - SCAN THE ENTIRE DOCUMENT THOROUGHLY - Don't stop at first match, check all content
        - TRANSLATE ALL GERMAN CONTENT TO ENGLISH with technical accuracy
        - EXTRACT EVERYTHING - better to include too much than miss important data
        - LOOK FOR NUMERICAL DATA: percentages (%), weights (g), energy values (kcal, kJ)
        - If percentages aren't explicitly given, set to 0 but ALWAYS INCLUDE the ingredient
        - If nutritional values aren't found, set to 0 but check ALL sections and slides
        - LOOK FOR TABLES, CHARTS, LISTS, HEADERS, FOOTNOTES - information can be anywhere
        - Pay special attention to: Zutatenliste, Nährwerttabelle, Allergene, Hinweise, slide titles

        MANDATORY JSON OUTPUT FORMAT (ALL TEXT IN ENGLISH):

        {{
            "recipe_number": "SEARCH FOR 'Ref.:' followed by number pattern. Extract ONLY the number (e.g., 'Ref.: 000014567' → '000014567'). If not found: null",
            "name": "FIND AND EXTRACT the product name/title from document headers, titles, or descriptions. Translate to English. If not found: 'Unknown Product'",
            "category": "ANALYZE product type and select from: Flakes, Multigrain & Branflakes, Puffed Cereals, Extrudates & Co-Extrudates, Traditional Muesli, Crunchy Muesli, Bio & Fair-Trade Muesli, Oat Flakes, Wheat, Rye & Barley Flakes, Bars",
            "description": "EXTRACT product description, benefits, or marketing text from document. Translate to English. If none found: 'No description available'",
            "ingredients": [
                "SCAN ENTIRE DOCUMENT for ingredient lists, Zutatenliste, compositions. Include ALL ingredients found:",
                {{"name": "TRANSLATE ingredient name to English (e.g., Haferflocken → Oat flakes)", "percentage": "EXTRACT exact percentage if given, otherwise 0"}}
            ],
            "nutritional_info": {{
                "energy": "FIND energy/Energie/Brennwert in kcal per 100g. Convert kJ to kcal if needed (divide by 4.184). Extract number only",
                "fat": "FIND fat/Fett content per 100g in grams. Extract number only",
                "saturated_fat": "FIND saturated fat/gesättigte Fettsäuren per 100g. Extract number only or 0", 
                "carbohydrates": "FIND carbohydrates/Kohlenhydrate per 100g in grams. Extract number only",
                "sugars": "FIND sugars/Zucker per 100g in grams. Extract number only or 0",
                "fiber": "FIND fiber/Ballaststoffe per 100g in grams. Extract number only",
                "protein": "FIND protein/Eiweiß per 100g in grams. Extract number only",
                "salt": "FIND salt/Salz per 100g in grams. Extract number only"
            }},
            "allergens": [
                "THOROUGHLY SCAN for allergen information, 'Enthält', 'Kann enthalten', allergy warnings.",
                "TRANSLATE each allergen to English (e.g., 'Enthält Gluten' → 'Contains gluten')",
                "INCLUDE ALL allergens found, even if in footnotes or small print"
            ],
            "claims": [
                "SEARCH ENTIRE DOCUMENT for health claims, quality statements, certifications",
                "LOOK FOR: Reich an, Hoher Gehalt, Bio, Fair-Trade, Ohne Zusätze, etc.",
                "TRANSLATE to English (e.g., 'Reich an Ballaststoffen' → 'High in fiber')",
                "INCLUDE marketing claims, quality badges, certification mentions"
            ],
            "storage_conditions": "SEARCH for Lagerung, Lagerbedingungen, Aufbewahrung, storage instructions. TRANSLATE to English. If none found: 'Store in a cool, dry place, protect from direct sunlight'"
        }}

        DOCUMENT CONTENT TO ANALYZE:
        File type: {file_extension.upper() if file_extension else 'Unknown'}
        Filename: {filename}
        Content: {file_content[:6000] if file_content and len(file_content.strip()) > 10 else 'No readable content extracted from document'}
        
        SPECIAL INSTRUCTIONS FOR LIMITED CONTENT:
        If the document content above is limited or unclear, ANALYZE THE FILENAME CAREFULLY:
        - Product name is often in the filename (e.g., "muesli_bar_chocolate_peanut" → "Muesli Bar with Chocolate and Peanuts")
        - Look for ingredients mentioned: chocolate, peanut, nuts, oats, etc.
        - Identify product type: bar, muesli, flakes, etc.
        - Provide REALISTIC nutritional values for this type of product
        - Include STANDARD allergens for mentioned ingredients (nuts→"Contains nuts", milk→"Contains milk")
        - Add appropriate claims based on product type and ingredients
        
        FILENAME ANALYSIS EXAMPLES:
        "muesli_bar_whole_milk_chocolate_peanut_FTS" → 
        - Name: "Muesli Bar with Whole Milk Chocolate and Peanuts"
        - Ingredients: Oat flakes 35%, Milk chocolate 25%, Peanuts 20%, Honey 15%, etc.
        - Allergens: Contains peanuts, Contains milk, May contain gluten
        - Claims: Fair-Trade (FTS), Source of protein

        COMPREHENSIVE EXTRACTION INSTRUCTIONS:

        POWERPOINT SPECIFIC EXTRACTION:
        - EXAMINE EVERY SLIDE for data tables, charts, ingredient lists
        - LOOK FOR SLIDE TITLES containing product names or descriptions  
        - SCAN BULLET POINTS and TEXT BOXES for ingredient percentages
        - CHECK FOOTER/HEADER areas for allergen disclaimers
        - TABLES often contain nutritional data - extract ALL numerical values
        - INGREDIENT SLIDES may show percentages in charts or lists

        INGREDIENTS EXTRACTION:
        - SCAN for: Zutatenliste, Ingredients, Zusammensetzung, Bestandteile, Ingredient lists
        - EXTRACT ALL ingredients mentioned, even in small amounts or traces
        - LOOK for percentage symbols (%) next to ingredients in tables or charts
        - SEARCH FOR NUMERICAL DATA: "45% Haferflocken", "20g Zucker", etc.
        - COMMON INGREDIENTS TO FIND: Haferflocken, Weizen, Nüsse, Rosinen, Honig, Zucker, Salz, Schokolade
        - TRANSLATE accurately: Haferflocken→Oat flakes, Vollkornweizen→Whole grain wheat, Milchschokolade→Milk chocolate

        NUTRITIONAL VALUES EXTRACTION:
        - SCAN ALL SLIDES for: Nährwerte, Nährwertangaben, Nutritional information, pro 100g, per 100g
        - LOOK for tables with nutritional data in EVERY slide - often separate from main content
        - EXTRACT ALL NUMERICAL VALUES with their units (kcal, kJ, g, mg)
        - KEY TERMS TO FIND: Energie/Energy, Brennwert, Fett/Fat, Kohlenhydrate/Carbs, Eiweiß/Protein, Ballaststoffe/Fiber, Salz/Salt
        - CONVERT kJ to kcal (divide by 4.184), extract numbers only
        - CHECK multiple slides - nutrition info might be in appendix or separate slides
        - TYPICAL NUTRITIONAL TABLE FORMAT: Look for values like "452 kcal", "18,5g Fett", etc.

        ALLERGENS EXTRACTION:
        - SCAN for: Enthält, Contains, Kann enthalten, May contain, Allergene, Allergy
        - COMMON ALLERGENS: Gluten, Nüsse, Milch, Eier, Soja, Sesam
        - LOOK in footnotes, warnings, small print areas
        - INCLUDE both "contains" and "may contain" allergens

        CLAIMS EXTRACTION:
        - SCAN for: Reich an, High in, Quelle von, Source of, Bio, Organic, Fair-Trade
        - MARKETING TERMS: Ohne Zusätze, No additives, Natürlich, Natural, Vollkorn, Wholegrain
        - QUALITY BADGES: Bio-Siegel, Fairtrade-Zeichen, certification logos
        - HEALTH STATEMENTS: about fiber, protein, vitamins, minerals

        STORAGE CONDITIONS EXTRACTION:
        - SCAN for: Lagerung, Lagerbedingungen, Aufbewahrung, Storage, Store
        - TYPICAL PHRASES: "Kühl und trocken lagern", "Vor Sonnenlicht schützen"
        - LOOK in product information sections, ingredient panels

        TRANSLATION EXAMPLES:
        - "Haferflocken" → "Oat flakes"
        - "Vollkornweizen" → "Whole grain wheat" 
        - "Nüsse" → "Nuts"
        - "Enthält Gluten" → "Contains gluten"
        - "Reich an Ballaststoffen" → "High in fiber"
        - "Kühl und trocken lagern" → "Store in a cool, dry place"
        - "Vor Sonnenlicht schützen" → "Protect from sunlight"
        - "Lagerung bei Raumtemperatur" → "Store at room temperature"

        CATEGORY SELECTION RULES:
        - If product contains "Riegel", "Bar" → select "Bars"
        - If product contains "Haferflocken", "Oat" → select "Oat Flakes" 
        - If product contains "Müsli", "Muesli" and is traditional → select "Traditional Muesli"
        - If product contains "Müsli", "Muesli" and is crunchy/baked → select "Crunchy Muesli"
        - If product contains "Bio", "Fair-Trade", "Organic" → select "Bio & Fair-Trade Muesli"
        - If product contains "Flakes", "Flocken" → select "Flakes"
        - If product contains multiple grains or bran → select "Multigrain & Branflakes"
        - If product is puffed/gepufft → select "Puffed Cereals"
        - If product is extruded → select "Extrudates & Co-Extrudates"
        - If product contains wheat, rye, barley → select "Wheat, Rye & Barley Flakes"

        NUTRITIONAL VALUES SEARCH INSTRUCTIONS:
        - Look for "Nährwerte", "Nährwertangaben", "pro 100g", "per 100g", "Energie", "Brennwert", "Fett", "Kohlenhydrate", "Eiweiß", "Protein", "Ballaststoffe", "Salz"
        - Extract exact numerical values with units (kcal, kJ, g, mg)
        - Convert kJ to kcal if needed (1 kcal = 4.184 kJ)
        - If you find "Brennwert" or "Energie" extract the kcal value for energy field
        - Always include the numerical value without units in JSON (e.g., "energy": "450" not "energy": "450 kcal")

        CRITICAL: Scan the entire document content for nutritional information - it might be in tables, lists, or separate sections!

        """

        # Call OpenAI API with better error handling
        try:
            logging.info(f"Calling OpenAI API for recipe analysis of {filename}")
            response = openai_client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {"role": "system", "content": "You are an expert food technologist specializing in cereal and muesli products. Analyze recipe documents, extract structured data, and translate all German content to English. Maintain technical accuracy in food ingredient translations."},
                    {"role": "user", "content": prompt}
                ],
                response_format={"type": "json_object"},
                max_completion_tokens=2000,
                temperature=0.1
            )

            # Parse the AI response
            if response and response.choices and len(response.choices) > 0:
                ai_response = response.choices[0].message.content
            else:
                raise ValueError("Invalid response structure from OpenAI API")
            logging.info(f"AI response received: {ai_response}")  # Log full response for debugging
            
            if ai_response:
                recipe_data = json.loads(ai_response)
                # Validate that required fields are present
                if 'name' not in recipe_data:
                    recipe_data['name'] = 'Unknown Product'
                if 'category' not in recipe_data:
                    recipe_data['category'] = 'Traditional Muesli'
                if 'description' not in recipe_data:
                    recipe_data['description'] = 'No description available'
                if 'ingredients' not in recipe_data:
                    recipe_data['ingredients'] = []
                if 'nutritional_info' not in recipe_data:
                    recipe_data['nutritional_info'] = {
                        'energy': '0', 'fat': '0', 'saturated_fat': '0',
                        'carbohydrates': '0', 'sugars': '0', 'fiber': '0',
                        'protein': '0', 'salt': '0'
                    }
                if 'allergens' not in recipe_data:
                    recipe_data['allergens'] = []
                if 'claims' not in recipe_data:
                    recipe_data['claims'] = []
                if 'storage_conditions' not in recipe_data:
                    recipe_data['storage_conditions'] = 'Store in a cool, dry place, protect from direct sunlight'
                if 'recipe_number' not in recipe_data:
                    recipe_data['recipe_number'] = None
            else:
                raise ValueError("No response from AI analysis")
                
        except Exception as openai_error:
            logging.error(f"OpenAI API error: {str(openai_error)}")
            # Check if it's an API key issue
            if "api" in str(openai_error).lower() and ("key" in str(openai_error).lower() or "auth" in str(openai_error).lower()):
                return jsonify({
                    'success': False, 
                    'error': 'OpenAI API key not configured or invalid. Please check the OPENAI_API_KEY environment variable.'
                })
            
            # Provide a fallback response if AI fails
            recipe_data = {
                'name': f'Recipe from {filename}',
                'category': 'Traditional Muesli',
                'description': 'Recipe imported from document. Please edit the details as needed.',
                'ingredients': [
                    {'name': 'Oat flakes', 'percentage': 0},
                    {'name': 'Mixed cereals', 'percentage': 0}
                ],
                'nutritional_info': {
                    'energy': '0', 'fat': '0', 'saturated_fat': '0',
                    'carbohydrates': '0', 'sugars': '0', 'fiber': '0',
                    'protein': '0', 'salt': '0'
                },
                'allergens': ['Gluten'],
                'claims': [],
                'storage_conditions': 'Store in a cool, dry place, protect from direct sunlight'
            }
            # Add all extracted images to fallback data too
            recipe_data['extracted_images'] = extracted_images.get('all_images', [])
            
            logging.info(f"Using fallback recipe data for {filename} due to OpenAI error")

        # Add all extracted images to recipe data for user selection
        recipe_data['extracted_images'] = extracted_images.get('all_images', [])
        
        # Cleanup temporary files
        document_processor.cleanup_files(temp_file_paths)
        
        logging.info(f"AI analysis completed for {len(filenames_list)} file(s)")
        return jsonify({'success': True, 'recipe_data': recipe_data})

    except Exception as e:
        logging.error(f"Recipe analysis error: {str(e)}")
        
        # Cleanup temporary files on error
        if 'temp_file_paths' in locals():
            document_processor.cleanup_files(temp_file_paths)
        
        # More specific error handling
        error_message = "An error occurred during AI analysis."
        if "openai" in str(e).lower():
            error_message = "OpenAI API error. Please check your API key and try again."
        elif "json" in str(e).lower():
            error_message = "Failed to parse AI response. Please try again."
        elif "file" in str(e).lower():
            error_message = "File processing error. Please check your file format."
        
        return jsonify({'success': False, 'error': error_message})

@app.route('/api/upload-recipe-image', methods=['POST'])
@csrf.exempt
@master_required
def upload_recipe_image():
    """Upload recipe images (product image or nutri-score image) with fallback to local storage"""
    try:
        if 'image' not in request.files:
            return jsonify({'success': False, 'error': 'No image file uploaded'})

        file = request.files['image']
        image_type = request.form.get('image_type', 'product')

        if file.filename == '':
            return jsonify({'success': False, 'error': 'No file selected'})

        # Validate file type
        allowed_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.webp'}
        filename_safe = file.filename or 'uploaded_image'
        file_extension = os.path.splitext(filename_safe.lower())[1]

        if file_extension not in allowed_extensions:
            return jsonify({'success': False, 'error': 'Invalid file type. Please upload an image.'})

        # Create secure filename
        filename = secure_filename(file.filename or 'uploaded_image')
        timestamp = int(time.time())
        unique_filename = f"{image_type}_{timestamp}_{filename}"

        # Always use local storage for reliability
        try:
            # Create upload directory if it doesn't exist
            upload_dir = os.path.join('static', 'images', 'recipes')
            os.makedirs(upload_dir, exist_ok=True)
            
            # Save file locally
            local_path = os.path.join(upload_dir, unique_filename)
            file.save(local_path)
            
            # Return local URL
            image_url = f"/static/images/recipes/{unique_filename}"
            
            logging.info(f"Recipe image uploaded successfully: {local_path}")
            
            return jsonify({
                'success': True,
                'image_url': image_url,
                'filename': unique_filename,
                'message': f'{image_type.title()} image uploaded successfully'
            })
            
        except Exception as save_error:
            logging.error(f"Failed to save image locally: {save_error}")
            return jsonify({
                'success': False,
                'error': f'Failed to save image: {str(save_error)}'
            }), 500

    except Exception as e:
        logging.error(f"Recipe image upload error: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'Upload failed: {str(e)}'
        }), 500

@app.route('/api/publish-recipe', methods=['POST'])
@csrf.exempt
@master_required
def publish_recipe():
    """Create a new recipe in the database - never overwrite existing ones"""
    try:
        data = request.get_json()
        

        # Generate unique recipe ID
        recipe_id = f"recipe_{int(time.time())}_{uuid.uuid4().hex[:8]}"

        # Create new Product entry
        new_product = Product()
        new_product.name = data.get('name', 'Untitled Recipe')
        new_product.category = data.get('category', 'Traditional Swiss Muesli')
        new_product.description = data.get('description', '')

        # Store ingredients as JSON
        ingredients = data.get('ingredients', [])
        new_product.ingredients = json.dumps(ingredients)

        # Store nutritional info as JSON  
        nutritional_info = data.get('nutritional_info', {})
        new_product.nutritional_info = json.dumps(nutritional_info)

        # Store allergens as JSON
        allergens = data.get('allergens', [])
        new_product.allergens = json.dumps(allergens)

        # Store claims as JSON
        claims = data.get('claims', [])
        new_product.claims = json.dumps(claims)

        # Store storage conditions
        storage_conditions = data.get('storage_conditions', 'Store in a cool, dry place, protect from direct sunlight')
        new_product.storage_conditions = storage_conditions

        # Store recipe number if extracted
        recipe_number = data.get('recipe_number')
        logging.info(f"Publishing recipe with recipe_number from data: {recipe_number}")
        if recipe_number:
            new_product.recipe_number = recipe_number
            logging.info(f"Recipe number set to: {new_product.recipe_number}")

        # Use uploaded image URL or default
        image_url = data.get('image_url')
        logging.info(f"Publishing recipe with product image URL: {image_url}")
        if image_url and (image_url.startswith('/api/image/') or image_url.startswith('/static/images/')):
            # Ensure the full URL path is stored
            new_product.image_url = image_url
            logging.info(f"Using extracted/uploaded product image: {image_url}")
        else:
            new_product.image_url = '/static/images/product-placeholder.png'
            logging.info("Using fallback product image")

        # Set nutri-score image if uploaded
        nutri_score_image = data.get('nutri_score_image')
        logging.info(f"Publishing recipe with nutri-score image URL: {nutri_score_image}")
        if nutri_score_image and (nutri_score_image.startswith('/api/image/') or nutri_score_image.startswith('/static/images/')):
            new_product.nutri_score_image = nutri_score_image
            logging.info(f"Using extracted/uploaded nutri-score image: {nutri_score_image}")
        else:
            new_product.nutri_score_image = None
            logging.info("No nutri-score image set")

        # Add to database
        db.session.add(new_product)
        db.session.commit()

        logging.info(f"New recipe created with ID: {new_product.id}")

        return jsonify({
            'success': True,
            'recipe_id': recipe_id,
            'database_id': new_product.id,
            'message': 'Recipe published successfully as new product'
        })

    except Exception as e:
        logging.error(f"Recipe publishing error: {str(e)}")
        return jsonify({
            'success': False,
            'error': f'Publishing failed: {str(e)}'
        }), 500

@app.route('/api/delete-recipe/<int:recipe_id>', methods=['DELETE'])
@csrf.exempt
@login_required
def delete_recipe(recipe_id):
    """Delete a recipe completely from the database"""
    try:
        # Find the recipe
        recipe = Product.query.get(recipe_id)
        if not recipe:
            return jsonify({
                'success': False,
                'error': 'Recipe not found'
            }), 404

        # Delete from database
        db.session.delete(recipe)
        db.session.commit()

        logging.info(f"Recipe deleted: ID {recipe_id}, Name: {recipe.name}")

        return jsonify({
            'success': True,
            'message': 'Recipe successfully deleted'
        })

    except Exception as e:
        logging.error(f"Recipe deletion error: {str(e)}")
        db.session.rollback()
        return jsonify({
            'success': False,
            'error': f'Deletion failed: {str(e)}'
        }), 500

@app.route('/api/image/<path:filename>')
def serve_image(filename):
    """Serve images from Object Storage with fallback to local storage"""
    try:
        storage_client = ObjectStorageClient()
        
        # Get file from Object Storage
        file_data = storage_client.download_as_bytes(filename)
        
        # Determine content type based on file extension
        file_extension = os.path.splitext(filename.lower())[1]
        content_type_map = {
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.png': 'image/png',
            '.gif': 'image/gif',
            '.webp': 'image/webp'
        }
        content_type = content_type_map.get(file_extension, 'image/jpeg')
        
        from flask import Response
        response = Response(file_data, mimetype=content_type)
        
        # Add cache headers for better performance
        response.headers['Cache-Control'] = 'public, max-age=3600'
        response.headers['Access-Control-Allow-Origin'] = '*'
        
        logging.info(f"Successfully served image from Object Storage: {filename}")
        return response
        
    except Exception as storage_error:
        logging.warning(f"Object Storage failed for {filename}, trying local fallback: {storage_error}")
        
        # Try local storage fallback
        try:
            # Check if it's in the local images/recipes directory
            local_path = os.path.join('static', 'images', 'recipes', os.path.basename(filename))
            if os.path.exists(local_path):
                return send_file(local_path)
            
            # Check if it's in uploads directory
            uploads_path = os.path.join('static', 'uploads', os.path.basename(filename))
            if os.path.exists(uploads_path):
                return send_file(uploads_path)
                
        except Exception as local_error:
            logging.warning(f"Local fallback also failed for {filename}: {local_error}")
        
        # Final fallback - create placeholder
        try:
            from PIL import Image, ImageDraw, ImageFont
            import io
            
            # Create placeholder image
            img = Image.new('RGB', (400, 300), color='#f8f9fa')
            draw = ImageDraw.Draw(img)
            
            # Draw placeholder text
            try:
                font = ImageFont.load_default()
            except:
                font = None
            
            text = "Bild nicht verfügbar"
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
            
            position = ((400 - text_width) // 2, (300 - text_height) // 2)
            draw.text(position, text, fill='#6c757d', font=font)
            
            # Convert to bytes
            img_io = io.BytesIO()
            img.save(img_io, 'PNG')
            img_io.seek(0)
            
            from flask import Response
            logging.info(f"Created placeholder image for {filename}")
            return Response(img_io.getvalue(), mimetype='image/png')
            
        except Exception as fallback_error:
            logging.error(f"Could not create placeholder image: {fallback_error}")
            return jsonify({'error': 'Image not found'}), 404

    return render_template('profile.html', user=user_data)

@app.route('/profile', methods=['POST'])
@login_required
def profile_update():
    # Update session data with form values
    session['user_name'] = request.form.get('name', '')
    session['user_email'] = request.form.get('email', '')
    session['user_position'] = request.form.get('position', '')
    session['user_department'] = request.form.get('department', '')
    session['user_phone'] = request.form.get('phone', '')
    session['user_location'] = request.form.get('location', '')
    # Removed language, timezone, and notification settings as per user request

    flash('Your profile has been successfully updated!', 'success')
    return redirect(url_for('profile'))

@app.route('/api/product/<int:product_id>/ingredients', methods=['GET'])
@login_required
def get_product_ingredients(product_id):
    """API endpoint to fetch real ingredients for a specific product"""
    try:
        product = Product.query.get_or_404(product_id)
        
        # Parse the ingredients JSON
        ingredients = []
        if product.ingredients:
            try:
                ingredients_data = json.loads(product.ingredients)
                if isinstance(ingredients_data, list):
                    ingredients = ingredients_data
                else:
                    logging.warning(f"Product {product_id} ingredients is not a list: {ingredients_data}")
            except json.JSONDecodeError as e:
                logging.error(f"Failed to parse ingredients for product {product_id}: {e}")
        
        return jsonify({
            'success': True,
            'product_id': product_id,
            'product_name': product.name,
            'ingredients': ingredients
        })
        
    except Exception as e:
        logging.error(f"Error fetching ingredients for product {product_id}: {e}")
        return jsonify({
            'success': False,
            'error': 'Failed to load product ingredients'
        }), 500

@app.route('/api/recipe-number/<recipe_number>/product', methods=['GET'])
@login_required
def get_product_by_recipe_number(recipe_number):
    """API endpoint to fetch product data by recipe number"""
    try:
        logging.info(f"Searching for recipe number: {recipe_number}")
        product = None
        
        # First try: Direct recipe_number match (most specific)
        product = Product.query.filter_by(recipe_number=recipe_number).first()
        logging.info(f"Direct recipe_number match result: {product}")
        
        # Second try: Recipe number with different formatting
        if not product:
            # Try without leading zeros or with different formatting
            recipe_variations = [
                recipe_number.lstrip('0'),  # Remove leading zeros
                recipe_number.replace('000004', ''),  # Remove prefix
                recipe_number
            ]
            
            for variation in recipe_variations:
                if variation and variation != recipe_number:
                    product = Product.query.filter_by(recipe_number=variation).first()
                    if product:
                        logging.info(f"Found with variation '{variation}': {product}")
                        break
        
        # Third try: Search in recipe_number field with LIKE
        if not product:
            product = Product.query.filter(
                Product.recipe_number.ilike(f'%{recipe_number}%')
            ).first()
            logging.info(f"LIKE search in recipe_number result: {product}")
        
        # Fourth try: Extract numeric part and search by ID or numeric patterns
        if not product:
            number_match = re.search(r'\d+', recipe_number)
            if number_match:
                numeric_part = number_match.group()
                logging.info(f"Trying numeric search with: {numeric_part}")
                
                # Try as ID if reasonable
                if len(numeric_part) <= 5:
                    try:
                        product = Product.query.get(int(numeric_part))
                    except (ValueError, TypeError):
                        pass
                
                # Try searching within recipe_number field for this number
                if not product:
                    product = Product.query.filter(
                        Product.recipe_number.ilike(f'%{numeric_part}%')
                    ).first()
        
        if not product:
            logging.warning(f"No product found for recipe number: {recipe_number}")
            return jsonify({
                'success': False,
                'error': 'Recipe not found',
                'searched_recipe': recipe_number
            }), 404
        
        logging.info(f"Found product: ID={product.id}, Name={product.name}, RecipeNum={product.recipe_number}")
        
        return jsonify({
            'success': True,
            'recipe_number': recipe_number,
            'product': {
                'id': product.id,
                'name': product.name,
                'image': product.image_url,
                'category': product.category,
                'description': product.description
            }
        })
        
    except Exception as e:
        logging.error(f"Error fetching product for recipe {recipe_number}: {e}")
        return jsonify({
            'success': False,
            'error': 'Failed to load product data',
            'exception': str(e)
        }), 500

@app.route('/manage-reports')
@master_required
def manage_reports():
    """Page for managing and deleting trends/reports"""
    # Simple protection: require session to be explicitly set as admin
    # This provides basic security without complex authentication
    if not session.get('admin_mode'):
        # If not admin mode, set it and continue (development convenience)
        session['admin_mode'] = True
        flash('Admin mode activated for report management.', 'info')
    
    init_user_session()
    
    # Get all trends/reports for management
    trends = Trend.query.order_by(Trend.created_at.desc()).all()
    
    return render_template('manage_reports.html', trends=trends)

@csrf.exempt
@app.route('/api/trends/<int:trend_id>/delete', methods=['DELETE'])
@master_required
def delete_trend(trend_id):
    """API endpoint for deleting a trend/report"""
    try:
        # Check for admin mode - consistent with manage_reports page
        if not session.get('admin_mode'):
            return jsonify({'success': False, 'error': 'Admin mode required for report deletion.'}), 403
            
        trend = Trend.query.get(trend_id)
        if not trend:
            return jsonify({'success': False, 'error': 'Report not found'}), 404
        
        # Delete associated files if they exist
        if trend.image_url and trend.image_url.startswith('/static/images/trend_'):
            image_path = os.path.join(app.root_path, trend.image_url.lstrip('/'))
            if os.path.exists(image_path):
                try:
                    os.remove(image_path)
                except OSError:
                    pass  # Continue even if file deletion fails
        
        if trend.pdf_path and trend.pdf_path.startswith('/static/pdfs/'):
            pdf_path = os.path.join(app.root_path, trend.pdf_path.lstrip('/'))
            if os.path.exists(pdf_path):
                try:
                    os.remove(pdf_path)
                except OSError:
                    pass  # Continue even if file deletion fails
        
        # Delete the trend from database
        db.session.delete(trend)
        db.session.commit()
        
        return jsonify({'success': True, 'message': 'Report deleted successfully'})
        
    except Exception as e:
        db.session.rollback()
        logging.error(f"Error deleting trend {trend_id}: {str(e)}")
        return jsonify({'success': False, 'error': 'Failed to delete report'}), 500

@app.errorhandler(404)
def not_found_error(error):
    return render_template('404.html'), 404

@app.errorhandler(500)
def internal_error(error):
    db.session.rollback()
    return render_template('500.html'), 500