import os
import tempfile
import zipfile
import base64
from pathlib import Path
import fitz  # PyMuPDF for PDF processing
from docx import Document  # python-docx for DOCX processing
from pptx import Presentation  # python-pptx for PowerPoint processing
from openpyxl import load_workbook  # openpyxl for Excel processing
import logging
from openai import OpenAI

logger = logging.getLogger(__name__)
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

class DocumentProcessor:
    """
    Process various document types and extract text content for AI analysis
    """
    
    def __init__(self):
        self.supported_types = [
            'application/pdf',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'text/plain',
            'image/jpeg',
            'image/jpg',
            'image/png'
        ]
    
    def is_supported(self, file_type):
        """Check if file type is supported"""
        return file_type in self.supported_types
    
    def process_file(self, file_path, file_type):
        """
        Process a file and extract text content based on file type
        
        Args:
            file_path (str): Path to the file
            file_type (str): MIME type of the file
            
        Returns:
            dict: Contains extracted text and metadata
        """
        try:
            if file_type == 'application/pdf':
                return self._process_pdf(file_path)
            elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return self._process_docx(file_path)
            elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return self._process_pptx(file_path)
            elif file_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                return self._process_xlsx(file_path)
            elif file_type == 'text/plain':
                return self._process_txt(file_path)
            elif file_type in ['image/jpeg', 'image/jpg', 'image/png']:
                return self._process_image(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
                
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            raise
    
    def _process_pdf(self, file_path):
        """Extract text from PDF file using PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            text_content = ""
            page_count = 0
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text_content += page.get_text("text")  # Updated PyMuPDF API
                page_count += 1
            
            doc.close()
            
            return {
                'text': text_content.strip(),
                'page_count': page_count,
                'word_count': len(text_content.split()),
                'file_type': 'PDF'
            }
            
        except Exception as e:
            raise Exception(f"Failed to process PDF: {str(e)}")
    
    def _process_docx(self, file_path):
        """Extract text from DOCX file using python-docx"""
        try:
            doc = Document(file_path)
            text_content = ""
            paragraph_count = 0
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
                    paragraph_count += 1
            
            # Extract text from tables if any
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_content += cell.text + " "
                    text_content += "\n"
            
            return {
                'text': text_content.strip(),
                'paragraph_count': paragraph_count,
                'word_count': len(text_content.split()),
                'file_type': 'Word Document'
            }
            
        except Exception as e:
            raise Exception(f"Failed to process DOCX: {str(e)}")
    
    def _process_txt(self, file_path):
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text_content = file.read()
            
            lines = text_content.split('\n')
            line_count = len([line for line in lines if line.strip()])
            
            return {
                'text': text_content.strip(),
                'line_count': line_count,
                'word_count': len(text_content.split()),
                'file_type': 'Text File'
            }
            
        except Exception as e:
            raise Exception(f"Failed to process TXT: {str(e)}")
    
    def _process_pptx(self, file_path):
        """Extract text from PowerPoint file using python-pptx"""
        try:
            prs = Presentation(file_path)
            text_content = ""
            slide_count = 0
            
            for slide in prs.slides:
                slide_count += 1
                # Extract text from shapes and tables
                for shape in slide.shapes:
                    if hasattr(shape, "has_table") and shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells]
                            if any(row_text):  # Only add non-empty rows
                                text_content += "\t".join(row_text) + "\n"
                    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        if shape.text and shape.text.strip():
                            text_content += shape.text + "\n"

                # Extract notes if available
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text and notes_text.strip():
                        text_content += f"[Notes: {notes_text}]\n"
                
                # Add slide separator
                text_content += "\n"
            
            return {
                'text': text_content.strip(),
                'slide_count': slide_count,
                'word_count': len(text_content.split()),
                'file_type': 'PowerPoint'
            }
            
        except Exception as e:
            raise Exception(f"Failed to process PowerPoint: {str(e)}")
    
    def _process_xlsx(self, file_path):
        """Extract text from Excel file using openpyxl"""
        try:
            workbook = load_workbook(file_path, data_only=True)
            text_content = ""
            sheet_count = len(workbook.sheetnames)
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content += f"--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    row_text = ' | '.join(row_data)
                    if row_text.strip():
                        text_content += row_text + "\n"
                
                text_content += "\n"
            
            return {
                'text': text_content.strip(),
                'sheet_count': sheet_count,
                'word_count': len(text_content.split()),
                'file_type': 'Excel'
            }
            
        except Exception as e:
            raise Exception(f"Failed to process Excel: {str(e)}")
    
    def _process_image(self, file_path):
        """Process image file using OpenAI Vision API"""
        try:
            # Read image and encode to base64
            with open(file_path, 'rb') as image_file:
                image_data = base64.b64encode(image_file.read()).decode('utf-8')
            
            # Determine image format from file extension
            file_ext = os.path.splitext(file_path)[1].lower()
            mime_type = 'image/jpeg' if file_ext in ['.jpg', '.jpeg'] else 'image/png'
            
            # Use OpenAI Vision API to analyze image
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": "Analyze this image and extract all text, product information, ingredients, nutritional data, or any recipe-related information visible. Provide a detailed description of what you see."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:{mime_type};base64,{image_data}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=1000
            )
            
            extracted_text = response.choices[0].message.content
            
            return {
                'text': extracted_text,
                'word_count': len(extracted_text.split()),
                'file_type': 'Image (Vision API)'
            }
            
        except Exception as e:
            logger.error(f"Vision API error: {str(e)}")
            # Fallback to placeholder if Vision API fails
            return {
                'text': '[IMAGE FILE - Could not analyze visual content]',
                'word_count': 0,
                'file_type': 'Image'
            }
    
    def process_multiple_files(self, files_data):
        """
        Process multiple files and combine their content
        
        Args:
            files_data (list): List of tuples (file_path, file_type, original_name)
            
        Returns:
            dict: Combined processing results
        """
        combined_text = ""
        file_summaries = []
        total_word_count = 0
        
        for file_path, file_type, original_name in files_data:
            try:
                result = self.process_file(file_path, file_type)
                
                # Add file separator
                combined_text += f"\n\n--- Content from {original_name} ---\n\n"
                combined_text += result['text']
                
                # Track file summary
                file_summaries.append({
                    'name': original_name,
                    'type': result['file_type'],
                    'word_count': result['word_count'],
                    'processed_successfully': True
                })
                
                total_word_count += result['word_count']
                
            except Exception as e:
                logger.error(f"Failed to process {original_name}: {str(e)}")
                file_summaries.append({
                    'name': original_name,
                    'type': 'Unknown',
                    'word_count': 0,
                    'processed_successfully': False,
                    'error': str(e)
                })
        
        return {
            'combined_text': combined_text.strip(),
            'file_summaries': file_summaries,
            'total_word_count': total_word_count,
            'total_files': len(files_data),
            'successful_files': len([f for f in file_summaries if f['processed_successfully']])
        }
    
    def save_uploaded_file(self, file_storage, upload_dir='/tmp/uploads'):
        """
        Save uploaded file to temporary directory
        
        Args:
            file_storage: Flask file storage object
            upload_dir: Directory to save files
            
        Returns:
            str: Path to saved file
        """
        try:
            # Create upload directory if it doesn't exist
            os.makedirs(upload_dir, exist_ok=True)
            
            # Generate unique filename
            filename = file_storage.filename
            temp_filename = f"{tempfile.gettempprefix()}{filename}"
            file_path = os.path.join(upload_dir, temp_filename)
            
            # Save file
            file_storage.save(file_path)
            
            return file_path
            
        except Exception as e:
            raise Exception(f"Failed to save uploaded file: {str(e)}")
    
    def cleanup_files(self, file_paths):
        """Clean up temporary files"""
        for file_path in file_paths:
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception as e:
                logger.warning(f"Failed to cleanup file {file_path}: {str(e)}")

# Global instance
document_processor = DocumentProcessor()